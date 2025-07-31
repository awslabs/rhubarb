# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: Apache-2.0
import base64
import logging
import warnings
import mimetypes
from io import BytesIO
from typing import Dict, List, Union, Optional

import boto3
import pdfplumber
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .image_validator import ImageValidator

logger = logging.getLogger(__name__)
warnings.filterwarnings("ignore", message="CropBox missing from /Page, defaulting to MediaBox")
warnings.filterwarnings("ignore", message="Glyph .* missing from font.*", category=UserWarning)


class FileConverter:
    def __init__(self, file_path: str, pages: List[int], s3_client: Optional[boto3.client], include_powerpoint_notes: bool = False):
        """
        Initialize the FileConverter object.

        Args:
            file_path (str): The path to the file (local or S3).
            pages (List[int]): List of page numbers to process.
            s3_client (boto3.client, optional): The boto3 S3 client. Defaults to None.
            include_powerpoint_notes (bool, optional): Whether to include PowerPoint slide notes. Defaults to False.
        """
        self.file_path = file_path
        self.s3_client = s3_client
        self.pages = pages
        self.include_powerpoint_notes = include_powerpoint_notes
        self.file_bytes, self.mime_type = self._get_file_bytes_and_mime_type()

    def _get_file_bytes_and_mime_type(self) -> tuple[bytes, str]:
        """
        Get the file bytes and MIME type of the file.

        Returns:
            tuple[bytes, str]: A tuple containing the file bytes and MIME type.
        """
        if self.file_path.startswith("s3://"):
            if self.s3_client is None:
                raise ValueError("S3 client is required for S3 file paths")
            bucket_name, key = self._parse_s3_path(self.file_path)
            response = self.s3_client.get_object(Bucket=bucket_name, Key=key)
            file_bytes = response["Body"].read()
        else:
            with open(self.file_path, "rb") as f:
                file_bytes = f.read()

        mime_type = self._get_mime_type(file_bytes)
        return file_bytes, mime_type

    def _get_mime_type(self, file_bytes: bytes) -> str:
        """
        Determine the MIME type of the file based on its contents.

        Args:
            file_bytes (bytes): The file bytes.

        Returns:
            str: The MIME type of the file.

        Raises:
            ValueError: If the file type is not supported.
        """
        mime_type, _ = mimetypes.guess_type(self.file_path)
        if mime_type is None:
            # Check for Office formats by file extension
            if self.file_path.lower().endswith(('.xlsx', '.xls')):
                mime_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            elif self.file_path.lower().endswith('.pptx'):
                mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            else:
                raise ValueError("Unsupported file type")
        return mime_type

    def _parse_s3_path(self, s3_path: str) -> tuple[str, str]:
        """
        Parse the S3 file path and extract the bucket name and key.

        Args:
            s3_path (str): The S3 file path.

        Returns:
            tuple[str, str]: A tuple containing the bucket name and key.
        """
        parts = s3_path[5:].partition("/")
        bucket_name = parts[0]
        key = parts[2]
        return bucket_name, key

    def convert_to_base64(self) -> List[Dict[str, Union[int, str]]]:
        """
        Convert the file to Base64 encoded string(s).

        Returns:
            List[Dict[str, Union[int, str]]]: A list of dictionaries containing the page number and Base64 encoded string.

        Raises:
            RuntimeError: If an error occurs during the file conversion process.
        """
        try:
            if self.mime_type == "image/jpeg" or self.mime_type == "image/png":
                if self.file_path.startswith("s3://"):
                    image_bytes = self.file_bytes
                else:
                    with open(self.file_path, "rb") as f:
                        image_bytes = f.read()
                validator = ImageValidator(image_bytes)
                validator.validate_image()
                base64_string = base64.b64encode(image_bytes).decode("utf-8")
                return [{"page": 1, "base64string": base64_string}]

            elif self.mime_type == "application/pdf":
                with pdfplumber.open(
                    BytesIO(self.file_bytes)
                    if self.file_path.startswith("s3://")
                    else self.file_path
                ) as pdf:
                    base64_strings = []
                    if self.pages == [0]:
                        page_nums = range(min(20, len(pdf.pages)))
                    else:
                        page_nums = [p - 1 for p in self.pages if p <= len(pdf.pages) and p > 0]

                    for page_num in page_nums:
                        page = pdf.pages[page_num]
                        img = page.to_image(resolution=150).original
                        img_bytes = BytesIO()
                        img.save(img_bytes, format="PNG")
                        base64_string = base64.b64encode(img_bytes.getvalue()).decode("utf-8")
                        base64_strings.append({"page": page_num + 1, "base64string": base64_string})
                return base64_strings

            elif self.mime_type == "image/tiff":
                with Image.open(
                    BytesIO(self.file_bytes)
                    if self.file_path.startswith("s3://")
                    else self.file_path
                ) as img:
                    base64_strings = []
                    if self.pages == [0]:
                        frame_nums = range(min(20, img.n_frames))
                    else:
                        frame_nums = [p - 1 for p in self.pages if p <= img.n_frames and p > 0]
                    for i in frame_nums:
                        img.seek(i)
                        img_byte_arr = BytesIO()
                        img.save(img_byte_arr, format="PNG")
                        base64_string = base64.b64encode(img_byte_arr.getvalue()).decode("utf-8")
                        base64_strings.append({"page": i + 1, "base64string": base64_string})
                return base64_strings
            elif self.mime_type in [
                "application/msword",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ]:
                try:
                    from docx import Document

                    document = Document(
                        BytesIO(self.file_bytes)
                        if self.file_path.startswith("s3://")
                        else self.file_path
                    )

                    base64_strings = []
                    page_count = len(
                        document.paragraphs
                    )  # Assuming paragraphs as a proxy for pages
                    if self.pages == [0]:
                        page_nums = range(min(20, page_count))
                    else:
                        page_nums = [p - 1 for p in self.pages if p <= page_count and p > 0]

                    for page_num in page_nums:
                        paragraph = document.paragraphs[page_num].text
                        img = Image.new(
                            "RGB", (800, 600), color=(255, 255, 255)
                        )  # Placeholder image for paragraph
                        d = ImageDraw.Draw(img)
                        d.text((10, 10), paragraph, fill=(0, 0, 0))
                        img_bytes = BytesIO()
                        img.save(img_bytes, format="PNG")
                        base64_string = base64.b64encode(img_bytes.getvalue()).decode("utf-8")
                        base64_strings.append({"page": page_num + 1, "base64string": base64_string})
                    return base64_strings
                except ImportError as e:
                    raise e
            elif self.mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                return self._convert_excel_to_base64()
            elif self.mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return self._convert_powerpoint_to_base64()
            else:
                logger.error("Unsupported file type")
                raise ValueError("Unsupported file type")
        except ImportError as e:
            logger.error(f"Error importing module: {str(e)}")
            raise ImportError(
                "The 'python-docx' library is not installed. Please install it to process .docx files.",
                "pip install python-docx",
            )
        except Exception as e:
            logger.error(f"Error converting file to Base64: {str(e)}")
            raise RuntimeError(f"Error converting file to Base64: {str(e)}")

    def convert_to_bytes(self) -> List[Dict[str, Union[int, bytes]]]:
        """
        Convert the file to bytes.

        Returns:
            List[Dict[str, Union[int, bytes]]]: A list of dictionaries containing the page number and bytes.

        Raises:
            RuntimeError: If an error occurs during the file conversion process.
        """
        try:
            if self.mime_type in ["image/jpeg", "image/png"]:
                if self.file_path.startswith("s3://"):
                    image_bytes = self.file_bytes
                else:
                    with open(self.file_path, "rb") as f:
                        image_bytes = f.read()
                validator = ImageValidator(image_bytes)
                validator.validate_image()
                return [{"page": 1, "image_bytes": image_bytes}]

            elif self.mime_type == "application/pdf":
                with pdfplumber.open(
                    BytesIO(self.file_bytes)
                    if self.file_path.startswith("s3://")
                    else self.file_path
                ) as pdf:
                    image_bytes_list = []
                    if self.pages == [0]:
                        page_nums = range(min(20, len(pdf.pages)))
                    else:
                        page_nums = [p - 1 for p in self.pages if p <= len(pdf.pages) and p > 0]

                    for page_num in page_nums:
                        page = pdf.pages[page_num]
                        img = page.to_image(resolution=150).original
                        img_byte_arr = BytesIO()
                        img.save(img_byte_arr, format="PNG")
                        image_bytes_list.append(
                            {"page": page_num + 1, "image_bytes": img_byte_arr.getvalue()}
                        )
                return image_bytes_list

            elif self.mime_type == "image/tiff":
                with Image.open(
                    BytesIO(self.file_bytes)
                    if self.file_path.startswith("s3://")
                    else self.file_path
                ) as img:
                    image_bytes_list = []
                    if self.pages == [0]:
                        frame_nums = range(min(20, img.n_frames))
                    else:
                        frame_nums = [p - 1 for p in self.pages if p <= img.n_frames and p > 0]
                    for i in frame_nums:
                        img.seek(i)
                        img_byte_arr = BytesIO()
                        img.save(img_byte_arr, format="PNG")
                        image_bytes_list.append(
                            {"page": i + 1, "image_bytes": img_byte_arr.getvalue()}
                        )
                return image_bytes_list

            elif self.mime_type in [
                "application/msword",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ]:
                try:
                    from docx import Document

                    document = Document(
                        BytesIO(self.file_bytes)
                        if self.file_path.startswith("s3://")
                        else self.file_path
                    )

                    image_bytes_list = []
                    page_count = len(
                        document.paragraphs
                    )  # Assuming paragraphs as a proxy for pages
                    if self.pages == [0]:
                        page_nums = range(min(20, page_count))
                    else:
                        page_nums = [p - 1 for p in self.pages if p <= page_count and p > 0]

                    for page_num in page_nums:
                        paragraph = document.paragraphs[page_num].text
                        img = Image.new(
                            "RGB", (800, 600), color=(255, 255, 255)
                        )  # Placeholder image for paragraph
                        d = ImageDraw.Draw(img)
                        d.text((10, 10), paragraph, fill=(0, 0, 0))
                        img_byte_arr = BytesIO()
                        img.save(img_byte_arr, format="PNG")
                        image_bytes_list.append(
                            {"page": page_num + 1, "image_bytes": img_byte_arr.getvalue()}
                        )
                    return image_bytes_list
                except ImportError as e:
                    raise e
            elif self.mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                return self._convert_excel_to_bytes()
            elif self.mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                return self._convert_powerpoint_to_bytes()
            else:
                logger.error("Unsupported file type")
                raise ValueError("Unsupported file type")
        except ImportError as e:
            logger.error(f"Error importing module: {str(e)}")
            raise ImportError(
                "The 'python-docx' library is not installed. Please install it to process .docx files.",
                "pip install python-docx",
            )
        except Exception as e:
            logger.error(f"Error converting file to bytes: {str(e)}")
            raise RuntimeError(f"Error converting file to bytes: {str(e)}")

    def _convert_excel_to_base64(self) -> List[Dict[str, Union[int, str]]]:
        """Convert Excel files to base64 encoded PNG images."""
        try:
            # Load workbook in read-only mode for memory efficiency
            if self.file_path.startswith("s3://"):
                workbook = load_workbook(BytesIO(self.file_bytes), read_only=True, data_only=True)
            else:
                workbook = load_workbook(self.file_path, read_only=True, data_only=True)
            
            base64_strings = []
            page_counter = 1
            
            # Process worksheets
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                
                # Get data dimensions
                max_row = worksheet.max_row
                max_col = worksheet.max_column
                
                if max_row == 1 and max_col == 1:  # Skip empty sheets
                    continue
                    
                # Handle large spreadsheets with chunking
                if max_row > 1000:  # Large sheet - use chunking
                    chunks = self._chunk_large_excel_sheet(worksheet, max_row, max_col)
                    for chunk_data, chunk_num in chunks:
                        if self.pages == [0] or page_counter in self.pages:
                            base64_string = self._render_excel_chunk_to_png(chunk_data, f"{sheet_name}_chunk_{chunk_num}")
                            base64_strings.append({"page": page_counter, "base64string": base64_string})
                        page_counter += 1
                        if len(base64_strings) >= 20:  # Respect 20-page limit
                            break
                else:  # Small sheet - render entirely
                    if self.pages == [0] or page_counter in self.pages:
                        # Extract all data
                        data = []
                        for row in worksheet.iter_rows(values_only=True, max_row=min(max_row, 500)):
                            data.append(row)
                        
                        base64_string = self._render_excel_data_to_png(data, sheet_name, max_col)
                        base64_strings.append({"page": page_counter, "base64string": base64_string})
                    page_counter += 1
                
                if len(base64_strings) >= 20:  # Global 20-page limit
                    break
            
            workbook.close()
            return base64_strings
            
        except Exception as e:
            logger.error(f"Error converting Excel file: {str(e)}")
            raise RuntimeError(f"Error converting Excel file: {str(e)}")

    def _convert_excel_to_bytes(self) -> List[Dict[str, Union[int, bytes]]]:
        """Convert Excel files to bytes."""
        base64_results = self._convert_excel_to_base64()
        bytes_results = []
        
        for result in base64_results:
            image_bytes = base64.b64decode(result["base64string"])
            bytes_results.append({"page": result["page"], "image_bytes": image_bytes})
        
        return bytes_results

    def _convert_powerpoint_to_base64(self) -> List[Dict[str, Union[int, str]]]:
        """Convert PowerPoint files to base64 encoded PNG images."""
        try:
            # Load presentation
            if self.file_path.startswith("s3://"):
                prs = Presentation(BytesIO(self.file_bytes))
            else:
                prs = Presentation(self.file_path)
            
            base64_strings = []
            
            # Determine which slides to process
            total_slides = len(prs.slides)
            if self.pages == [0]:
                slide_nums = range(min(20, total_slides))
            else:
                slide_nums = [p - 1 for p in self.pages if p <= total_slides and p > 0]
            
            for slide_idx in slide_nums:
                slide = prs.slides[slide_idx]
                base64_string = self._render_slide_to_png(slide, slide_idx + 1, self.include_powerpoint_notes)
                base64_strings.append({"page": slide_idx + 1, "base64string": base64_string})
            
            return base64_strings
            
        except Exception as e:
            logger.error(f"Error converting PowerPoint file: {str(e)}")
            raise RuntimeError(f"Error converting PowerPoint file: {str(e)}")

    def _convert_powerpoint_to_bytes(self) -> List[Dict[str, Union[int, bytes]]]:
        """Convert PowerPoint files to bytes."""
        base64_results = self._convert_powerpoint_to_base64()
        bytes_results = []
        
        for result in base64_results:
            image_bytes = base64.b64decode(result["base64string"])
            bytes_results.append({"page": result["page"], "image_bytes": image_bytes})
        
        return bytes_results

    def _chunk_large_excel_sheet(self, worksheet, max_row, max_col, chunk_size=100):
        """Break large Excel sheets into manageable chunks."""
        chunks = []
        
        # Always include headers (first 5 rows)
        header_data = []
        for row in worksheet.iter_rows(min_row=1, max_row=min(5, max_row), values_only=True):
            header_data.append(row)
        
        chunks.append((header_data, "headers"))
        
        # Process data in chunks
        for start_row in range(6, max_row + 1, chunk_size):
            end_row = min(start_row + chunk_size - 1, max_row)
            chunk_data = []
            
            # Add header context to each chunk
            chunk_data.extend(header_data)
            chunk_data.append(["..."] * max_col)  # Separator
            
            # Add chunk data
            for row in worksheet.iter_rows(min_row=start_row, max_row=end_row, values_only=True):
                chunk_data.append(row)
                
            chunks.append((chunk_data, f"rows_{start_row}_{end_row}"))
            
            # Limit chunks to prevent excessive processing
            if len(chunks) >= 10:
                break
        
        return chunks

    def _render_excel_data_to_png(self, data, sheet_name, max_col):
        """Render Excel data as a formatted PNG image."""
        # Set matplotlib to use a non-interactive backend
        plt.switch_backend('Agg')
        
        fig, ax = plt.subplots(figsize=(max(12, max_col * 1.2), max(8, len(data) * 0.3)))
        ax.axis('off')
        
        # Create table
        table_data = []
        for row in data:
            formatted_row = []
            for cell in row:
                if cell is None:
                    formatted_row.append("")
                else:
                    formatted_row.append(str(cell)[:50])  # Limit cell content
            table_data.append(formatted_row)
        
        # Limit columns for display
        if max_col > 15:
            table_data = [row[:15] + ["..."] for row in table_data]
        
        table = ax.table(cellText=table_data, loc='center', cellLoc='left')
        table.auto_set_font_size(False)
        table.set_fontsize(8)
        table.scale(1, 1.5)
        
        # Style header row
        for i in range(min(len(table_data[0]), 16)):
            table[(0, i)].set_facecolor('#4CAF50')
            table[(0, i)].set_text_props(weight='bold', color='white')
        
        plt.title(f"Sheet: {sheet_name}", fontsize=14, fontweight='bold')
        
        # Convert to PNG bytes
        img_bytes = BytesIO()
        plt.savefig(img_bytes, format='PNG', dpi=150, bbox_inches='tight')
        plt.close()
        
        return base64.b64encode(img_bytes.getvalue()).decode('utf-8')

    def _render_excel_chunk_to_png(self, chunk_data, chunk_name):
        """Render Excel chunk data as a formatted PNG image."""
        return self._render_excel_data_to_png(chunk_data, chunk_name, len(chunk_data[0]) if chunk_data else 1)

    def _render_slide_to_png(self, slide, slide_number, include_notes=False):
        """Render PowerPoint slide as PNG image."""
        # Set matplotlib to use a non-interactive backend
        plt.switch_backend('Agg')
        
        # Adjust figure height if including notes
        fig_height = 12 if include_notes else 9
        fig, ax = plt.subplots(figsize=(16, fig_height))  # 16:9 or 16:12 aspect ratio
        ax.set_xlim(0, 10)
        ax.set_ylim(0, fig_height * 0.83)  # Adjust y-limit based on figure height
        ax.axis('off')
        
        # Set white background
        fig.patch.set_facecolor('white')
        
        y_position = fig_height * 0.78
        
        # Extract and render slide content
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                # Text content
                text = shape.text.strip()[:200]  # Limit text length
                
                # Determine text size based on shape properties
                font_size = 12
                if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs and para.runs[0].font.size:
                        font_size = max(8, min(24, para.runs[0].font.size.pt))
                
                ax.text(0.5, y_position, text, fontsize=font_size, 
                       ha='left', va='top', wrap=True, 
                       bbox=dict(boxstyle="round,pad=0.3", facecolor="lightblue", alpha=0.7))
                y_position -= 0.8
                
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # pylint: disable=no-member
                # Image placeholder
                ax.add_patch(patches.Rectangle((0.5, y_position-0.5), 3, 1, 
                                             linewidth=1, edgecolor='gray',
                                             facecolor='lightgray', alpha=0.5))
                ax.text(2, y_position, "[IMAGE]", ha='center', va='center', fontsize=10)
                y_position -= 1.2
            
            # Stop if we're running out of space (but leave room for notes if needed)
            min_y_position = 3 if include_notes else 1
            if y_position < min_y_position:
                break
        
        # Add slide notes if requested and available
        if include_notes and slide.has_notes_slide:
            try:
                notes_slide = slide.notes_slide
                notes_text = ""
                
                # Extract notes text from notes slide
                for shape in notes_slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        # Skip the slide thumbnail placeholder
                        if "Click to add notes" not in shape.text:
                            notes_text += shape.text.strip() + " "
                
                if notes_text.strip():
                    # Add separator line
                    ax.axhline(y=2.5, xmin=0.05, xmax=0.95, color='gray', linestyle='--', linewidth=1)
                    
                    # Add notes section title
                    ax.text(0.5, 2.3, "Speaker Notes:", fontsize=12, fontweight='bold', 
                           ha='left', va='top', color='darkblue')
                    
                    # Add notes text (limit length)
                    notes_text = notes_text.strip()[:400]  # Limit notes length
                    ax.text(0.5, 1.9, notes_text, fontsize=10, 
                           ha='left', va='top', wrap=True, 
                           bbox=dict(boxstyle="round,pad=0.3", facecolor="lightyellow", alpha=0.7))
                           
            except Exception as e:
                # If notes extraction fails, just log and continue
                logger.warning(f"Could not extract notes from slide {slide_number}: {str(e)}")
        
        plt.title(f"Slide {slide_number}", fontsize=16, fontweight='bold', pad=20)
        
        # Convert to PNG bytes
        img_bytes = BytesIO()
        plt.savefig(img_bytes, format='PNG', dpi=150, bbox_inches='tight', facecolor='white')
        plt.close()
        
        return base64.b64encode(img_bytes.getvalue()).decode('utf-8')
